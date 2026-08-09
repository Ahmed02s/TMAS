import random
import re
import json
from pathlib import Path
import httpx
from datetime import datetime, timedelta, timezone
from typing import Any
import re as _re
from postgrest import exceptions as _postgrest_exceptions

from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel, Field

from app.core.security import get_current_claims, get_current_claims_optional, require_roles
from app.core.supabase_client import ensure_supabase_enabled, supabase, supabase_failed, supabase_error_message
from app.core.config import QROK_API_KEY, QROK_API_URL

BASE_DIR = Path(__file__).resolve().parent.parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"

# Small allowance for network/client latency when enforcing a student's per-attempt time limit.
TIME_LIMIT_GRACE_SECONDS = 20


def _verify_acting_as_self(claims: dict[str, Any], student_id: str) -> None:
    """Without this, any logged-in account could start/submit/forfeit a quiz attempt on
    behalf of an arbitrary student_id in the request body — the endpoint had no way to know
    who was really calling it. Admins/lecturers aren't exempted: nothing about starting or
    submitting a student's quiz attempt is a legitimate lecturer/admin action."""
    if str(claims.get('sub') or '') != str(student_id or ''):
        raise HTTPException(status_code=403, detail='You can only do this for your own account')


def _clean_material_title(title: str) -> str:
    title = re.sub(r'\.[^.]+$', '', title)
    title = title.replace('_', ' ').replace('-', ' ')
    title = re.sub(r'^\s*\d+[\s._-]*', '', title)
    title = re.sub(r'\b(lecture|lectures|notes|note|slide|slides|pptx|ppt|pdf|docx|doc)\b', '', title, flags=re.IGNORECASE)
    title = re.sub(r'\s+', ' ', title).strip()
    title = re.sub(r'\s*[:\-\(\)]\s*', ' ', title).strip()
    if not title or re.fullmatch(r'\d+', title):
        return 'course concept'
    title = title.title()
    return title

router = APIRouter(prefix='/api/quizzes', tags=['quizzes'])


def _parse_datetime(value: str | datetime | None) -> datetime | None:
    if not value:
        return None
    if isinstance(value, datetime):
        dt = value
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        return dt

    cleaned = value.strip()
    if cleaned.endswith('Z'):
        cleaned = cleaned[:-1] + '+00:00'

    # datetime.fromisoformat() is strict about fractional-second precision on Python < 3.11
    # (it rejects anything other than exactly 6 digits). JS's Date.toISOString() always emits
    # exactly 3 digits (milliseconds), so normalize any fractional-seconds component to 6
    # digits here — otherwise a perfectly valid timestamp silently fails to parse and open/close
    # dates end up as None (quiz stuck "scheduled"/locked forever), independent of Python version.
    frac_match = re.match(r'^(.*T\d{2}:\d{2}:\d{2})\.(\d+)([+-]\d{2}:\d{2}|)$', cleaned)
    if frac_match:
        frac = frac_match.group(2)[:6].ljust(6, '0')
        cleaned = f"{frac_match.group(1)}.{frac}{frac_match.group(3)}"

    try:
        dt = datetime.fromisoformat(cleaned)
        # Always ensure UTC-aware so comparisons with datetime.now(UTC) are safe
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        return dt
    except ValueError:
        for fmt in (
            '%Y-%m-%d %H:%M',
            '%Y-%m-%dT%H:%M',
            '%Y-%m-%d %H:%M:%S',
            '%Y-%m-%dT%H:%M:%S',
            '%Y-%m-%d %H:%M:%S%z',
            '%Y-%m-%dT%H:%M:%S%z',
            '%Y-%m-%dT%H:%M:%S.%f%z',
            '%Y-%m-%dT%H:%M:%S.%f',
        ):
            try:
                dt = datetime.strptime(cleaned, fmt)
                if dt.tzinfo is None:
                    dt = dt.replace(tzinfo=timezone.utc)
                return dt
            except ValueError:
                continue
        return None


def _format_datetime(value: datetime) -> str:
    formatted = value.replace(microsecond=0).isoformat()
    if formatted.endswith('+00:00'):
        formatted = formatted[:-6] + 'Z'
    return formatted


def _grade_from_score(score: int) -> str:
    if score >= 90:
        return 'A'
    if score >= 80:
        return 'B'
    if score >= 70:
        return 'C'
    if score >= 60:
        return 'D'
    return 'F'


def _build_question_seed(course: str, material_titles: list[str], difficulty: str, index: int) -> int:
    base = f"{course}|{'|'.join(material_titles)}|{difficulty}|{index}"
    return abs(hash(base)) % (2**32)


def _normalize_tier(tier: str) -> str:
    tier = (tier or '').strip().title()
    if tier in {'Foundational', 'Intermediate', 'Mastery'}:
        return tier
    return 'Foundational'


def _tier_default_difficulty(tier: str) -> str:
    if tier == 'Foundational':
        return 'Easy'
    if tier == 'Intermediate':
        return 'Medium'
    if tier == 'Mastery':
        return 'Hard'
    return 'Mixed'


def _quiz_is_open(quiz: dict[str, Any], now: datetime) -> bool:
    open_dt = _parse_datetime(quiz.get('open_date'))
    close_dt = _parse_datetime(quiz.get('close_date'))
    if open_dt and now < open_dt:
        return False
    if close_dt and now > close_dt:
        return False
    return True


def _normalize_difficulty(difficulty: str, rng: random.Random) -> str:
    difficulty = difficulty.strip().capitalize()
    if difficulty in {'Easy', 'Medium', 'Hard'}:
        return difficulty
    return rng.choice(['Easy', 'Medium', 'Hard'])


def _extract_text_from_docx(path: str) -> str:
    try:
        import docx
        doc = docx.Document(path)
        text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        if text.strip():
            return text
    except Exception:
        pass
    try:
        import zipfile
        lines = []
        with zipfile.ZipFile(path, 'r') as z:
            if 'word/document.xml' in z.namelist():
                xml_content = z.read('word/document.xml').decode('utf-8', errors='ignore')
                texts = re.findall(r'<w:t[^>]*>(.*?)</w:t>', xml_content)
                if texts:
                    lines.extend(t.strip() for t in texts if t.strip())
        return '\n'.join(lines)
    except Exception:
        return ''


def _extract_text_from_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        presentation = Presentation(path)
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text:
                        lines.append(text)
        if lines:
            return '\n'.join(lines)
    except Exception:
        pass
    try:
        import zipfile
        lines = []
        with zipfile.ZipFile(path, 'r') as z:
            for name in sorted(z.namelist()):
                if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                    xml_content = z.read(name).decode('utf-8', errors='ignore')
                    texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', xml_content)
                    if texts:
                        slide_num = re.search(r'slide(\d+)\.xml', name)
                        lines.append(f"--- Slide {slide_num.group(1) if slide_num else ''} ---")
                        lines.extend(t.strip() for t in texts if t.strip())
        return '\n'.join(lines)
    except Exception:
        return ''


def _extract_text_from_pdf(path: str) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(path)
        text = '\n'.join(page.extract_text() for page in reader.pages if page.extract_text())
        if text.strip():
            return text
    except Exception:
        pass
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(path)
        text = '\n'.join(page.extract_text() for page in reader.pages if page.extract_text())
        if text.strip():
            return text
    except Exception:
        pass
    try:
        content = Path(path).read_bytes().decode('utf-8', errors='ignore')
        matches = re.findall(r'\((.*?)\)\s*Tj|\((.*?)\)\s*TJ', content)
        lines = [m[0] or m[1] for m in matches if (m[0] or m[1]).strip()]
        return '\n'.join(lines)
    except Exception:
        return ''


def _extract_text_from_file(path: str) -> str:
    if not path:
        return ''
    file = Path(path)
    if not file.exists():
        return ''
    suffix = file.suffix.lower()
    try:
        if suffix in {'.txt', '.md', '.py', '.json', '.csv'}:
            return file.read_text(encoding='utf-8', errors='ignore')
        if suffix in {'.docx', '.doc'}:
            return _extract_text_from_docx(str(file))
        if suffix in {'.pptx', '.ppt'}:
            return _extract_text_from_pptx(str(file))
        if suffix == '.pdf':
            return _extract_text_from_pdf(str(file))
    except Exception:
        return ''
    return ''


def _summarize_material_text(text: str, fallback: str) -> str:
    if not text or len(text.strip()) < 50:
        return _clean_material_title(fallback)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return _clean_material_title(fallback)
    summary = lines[0]
    if len(summary) < 40 and len(lines) > 1:
        summary = summary + ' ' + lines[1][:120]
    summary = re.sub(r'\s+', ' ', summary).strip()
    return summary[:400]


def _infer_topic_context(topic: str, course: str) -> dict[str, str]:
    normalized = topic.strip().lower()
    descriptor = topic
    if 'specification' in normalized or 'requirements' in normalized:
        return {
            'descriptor': 'the concept of specification and requirements',
            'correct': f'defining the expected behavior and constraints of a system in {course}',
            'wrong1': f'a mathematical proof technique unrelated to {course}',
            'wrong2': f'a design aesthetic topic not central to {course}',
            'wrong3': f'a historical overview unrelated to the course objectives',
            'fill_blank': 'requirements',
            'short': f'It explains how to document what a system should do in {course}.',
            'true_statement': f"In {course}, '{topic}' most likely refers to defining what a system must do.",
        }
    if 'design' in normalized or 'architecture' in normalized:
        return {
            'descriptor': 'the concept of software design',
            'correct': f'planning the structure and organization of software for {course}',
            'wrong1': f'writing documentation unrelated to software structure',
            'wrong2': f'creating artwork rather than technical solutions',
            'wrong3': f'covering only hardware installation details',
            'fill_blank': 'design',
            'short': f'It explains how to organize software components to meet {course} goals.',
            'true_statement': f"In {course}, '{topic}' most likely refers to planning the structure of a software system.",
        }
    if 'algorithm' in normalized or 'complexity' in normalized:
        return {
            'descriptor': 'the concept of algorithms and analysis',
            'correct': f'defining efficient problem-solving steps for {course}',
            'wrong1': f'comparing unrelated business strategies',
            'wrong2': f'describing artistic techniques outside of {course}',
            'wrong3': f'listing administrative procedures rather than computation',
            'fill_blank': 'algorithm',
            'short': f'It explains how to solve problems efficiently in {course}.',
            'true_statement': f"In {course}, '{topic}' most likely refers to an algorithmic process or analysis.",
        }
    if 'data structure' in normalized or 'data' in normalized or 'tree' in normalized or 'graph' in normalized:
        return {
            'descriptor': 'the concept of data organization',
            'correct': f'structuring data so it can be accessed and used efficiently in {course}',
            'wrong1': f'describing unrelated physical storage hardware',
            'wrong2': f'explaining only aesthetic presentation details',
            'wrong3': f'covering only the history of computing without practical use',
            'fill_blank': 'structure',
            'short': f'It explains how to organize and access data effectively in {course}.',
            'true_statement': f"In {course}, '{topic}' most likely refers to organizing data for efficient use.",
        }
    return {
        'descriptor': f'the concept of {topic}',
        'correct': f'clarifying how {topic} supports the learning goals of {course}',
        'wrong1': f'focusing on unrelated historical details outside {course}',
        'wrong2': f'describing design choices unrelated to {course} outcomes',
        'wrong3': f'presenting a topic that does not contribute to practical understanding in {course}',
        'fill_blank': topic.lower(),
        'short': f'It explains how {topic} is used to meet the course objectives in {course}.',
        'true_statement': f"In {course}, '{topic}' most likely refers to an important course concept.",
    }


def _generate_fallback_question(course: str, material_titles: list[str], difficulty: str, question_type: str, index: int) -> dict[str, Any]:
    topic = (material_titles or [course])[index % (len(material_titles) or 1)]
    topic_ctx = _infer_topic_context(topic, course)
    marks_map = {'Easy': 1, 'Medium': 2, 'Hard': 3}
    marks = marks_map.get(difficulty, 2)

    if question_type == 'True/False':
        return {
            'id': index + 1,
            'type': 'True/False',
            'difficulty': difficulty,
            'topic': topic,
            'marks': marks,
            'question': topic_ctx['true_statement'],
            'options': ['True', 'False'],
            'answer': 'True',
            'explanation': topic_ctx['short'],
        }
    elif question_type == 'Fill in the Blank':
        return {
            'id': index + 1,
            'type': 'Fill in the Blank',
            'difficulty': difficulty,
            'topic': topic,
            'marks': marks,
            'question': f"In {course}, the core focus of '{topic}' relates to _____.",
            'options': [],
            'answer': topic_ctx['fill_blank'],
            'explanation': topic_ctx['short'],
        }
    elif question_type == 'Short Answer':
        return {
            'id': index + 1,
            'type': 'Short Answer',
            'difficulty': difficulty,
            'topic': topic,
            'marks': marks,
            'question': f"Briefly explain the primary role of '{topic}' in the context of {course}.",
            'options': [],
            'answer': topic_ctx['correct'],
            'explanation': topic_ctx['short'],
        }
    else:  # MCQ
        options = [
            topic_ctx['correct'],
            topic_ctx['wrong1'],
            topic_ctx['wrong2'],
            topic_ctx['wrong3'],
        ]
        return {
            'id': index + 1,
            'type': 'MCQ',
            'difficulty': difficulty,
            'topic': topic,
            'marks': marks,
            'question': f"Which statement best describes '{topic}' in {course}?",
            'options': options,
            'answer': topic_ctx['correct'],
            'explanation': topic_ctx['short'],
        }


def _generate_question(course: str, material_titles: list[str], difficulty: str, question_type: str, index: int) -> dict[str, Any]:
    llm_q = _call_llm_for_questions(course, material_titles, question_type, difficulty, index)
    if not llm_q:
        llm_q = _generate_fallback_question(course, material_titles, difficulty, question_type, index)

    llm_q['id'] = index + 1
    llm_q['type'] = question_type
    llm_q['difficulty'] = llm_q.get('difficulty', difficulty)
    llm_q['marks'] = int(llm_q.get('marks', 0) or 0)
    return llm_q


def _call_llm_for_questions(course: str, material_titles: list[str], question_type: str, difficulty: str, index: int) -> dict[str, Any] | None:
    prompt = (
        f"Generate a single {question_type} question for the course '{course}' at {difficulty} difficulty.\n"
        f"Base the question on these materials: {', '.join(material_titles)}.\n"
        "Return ONLY a JSON object with keys:\n"
        "- question: string\n"
        "- options: array of strings (for MCQ specify 4 options; for True/False specify ['True', 'False']; for Short Answer/Fill in the Blank leave as empty array [])\n"
        "- answer: string (matching exact correct option text for MCQ/True/False, or answer text for short answer)\n"
        "- explanation: string\n"
        "- difficulty: string ('Easy', 'Medium', or 'Hard')\n"
        "- marks: integer (1 for Easy, 2 for Medium, 3 for Hard)\n"
        "Do not output markdown codeblocks or extra text."
    )

    from app.core.config import GEMINI_API_KEY, OPENAI_API_KEY

    raw_text = None

    # Step 1: Try Primary (Groq / Qrok)
    if QROK_API_KEY:
        try:
            api_url = QROK_API_URL or 'https://api.groq.com/openai/v1/chat/completions'
            headers = {"Authorization": f"Bearer {QROK_API_KEY}", "Content-Type": "application/json"}
            body = {
                "model": "llama-3.3-70b-versatile",
                "messages": [
                    {"role": "system", "content": "You are a professional university professor creating high quality exam questions. Output strictly valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 512,
                "response_format": {"type": "json_object"}
            }
            with httpx.Client(timeout=12.0) as client:
                resp = client.post(api_url, headers=headers, json=body)
                if resp.status_code == 200:
                    raw_text = resp.json()['choices'][0]['message']['content']
        except Exception:
            pass

    # Step 2: Try Secondary (Google Gemini 1.5 Flash)
    if not raw_text and GEMINI_API_KEY:
        try:
            gemini_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={GEMINI_API_KEY}"
            g_body = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"response_mime_type": "application/json"}
            }
            with httpx.Client(timeout=12.0) as client:
                g_resp = client.post(gemini_url, json=g_body)
                if g_resp.status_code == 200:
                    raw_text = g_resp.json()['candidates'][0]['content']['parts'][0]['text']
        except Exception:
            pass

    # Step 3: Try Tertiary (OpenAI gpt-4o-mini)
    if not raw_text and OPENAI_API_KEY:
        try:
            o_url = "https://api.openai.com/v1/chat/completions"
            o_headers = {"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"}
            o_body = {
                "model": "gpt-4o-mini",
                "messages": [{"role": "user", "content": prompt}],
                "response_format": {"type": "json_object"}
            }
            with httpx.Client(timeout=12.0) as client:
                o_resp = client.post(o_url, headers=o_headers, json=o_body)
                if o_resp.status_code == 200:
                    raw_text = o_resp.json()['choices'][0]['message']['content']
        except Exception:
            pass

    if not raw_text:
        return None

    try:
        json_text = json.loads(raw_text)
    except Exception:
        m = re.search(r"\{.*\}", raw_text, re.S)
        json_text = json.loads(m.group(0)) if m else None

    if not isinstance(json_text, dict):
        return None

    question_text = json_text.get('question', '').strip()
    if not question_text:
        return None

    options = json_text.get('options', []) or []
    if not isinstance(options, list):
        options = []
    options = [str(opt).strip() for opt in options if opt is not None]

    raw_ans = json_text.get('answer', '')
    answer = str(raw_ans).strip() if raw_ans is not None else ''

    if options and answer.isdigit():
        idx = int(answer)
        if 0 <= idx < len(options):
            answer = options[idx]
        elif 1 <= idx <= len(options):
            answer = options[idx - 1]

    explanation = str(json_text.get('explanation', '') or '').strip()
    normalized_difficulty = json_text.get('difficulty', difficulty)
    marks = int(json_text.get('marks', {'Easy': 1, 'Medium': 2, 'Hard': 3}.get(normalized_difficulty, 2)))

    return {
        'id': index + 1,
        'type': question_type,
        'difficulty': normalized_difficulty,
        'topic': (material_titles or [course])[index % (len(material_titles) or 1)],
        'marks': marks,
        'question': question_text,
        'options': options,
        'answer': answer,
        'explanation': explanation or f"Base answer on {(material_titles or [course])[0]} in context of {course}.",
    }



def _is_missing_table_error(exc: Exception) -> bool:
    message = str(exc)
    if 'column of' in message:
        return False
    return 'Could not find the table' in message or ("Could not find the '" in message and "' table" in message) or ('relation "' in message and 'does not exist' in message) or ('relation' in message and 'does not exist' in message)


def _finalize_overdue_attempts(quizzes: list[dict[str, Any]], student_id: str) -> list[dict[str, Any]]:
    now = datetime.now(timezone.utc)
    overdue_quiz_ids = []
    for quiz in quizzes:
        due_dt = _parse_datetime(quiz.get('due_date'))
        if due_dt and due_dt < now:
            overdue_quiz_ids.append(quiz.get('id'))

    if not overdue_quiz_ids:
        return quizzes

    try:
        response = supabase.table('quiz_attempts').select('*').eq('student_id', student_id).in_('quiz_id', [qid for qid in overdue_quiz_ids if qid is not None]).execute()
    except _postgrest_exceptions.APIError as exc:
        if _is_missing_table_error(exc):
            return quizzes
        raise

    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase fetch quiz attempts failed'))

    all_existing = response.data or []
    finalized_quiz_ids = {a['quiz_id'] for a in all_existing if a.get('status') in ('completed', 'missed')}
    in_progress_by_quiz = {a['quiz_id']: a for a in all_existing if a.get('status') == 'in_progress'}
    existing_any = {a['quiz_id'] for a in all_existing}

    # Update any in-progress attempts for expired quizzes to 'missed'
    for quiz_id, attempt in in_progress_by_quiz.items():
        if quiz_id in overdue_quiz_ids and quiz_id not in finalized_quiz_ids:
            supabase.table('quiz_attempts').update({
                'status': 'missed',
                'attempted_at': datetime.now(timezone.utc).isoformat(),
            }).eq('id', attempt['id']).execute()

    # Insert missed records for overdue quizzes with no attempt at all
    missing = [quiz for quiz in quizzes if quiz.get('id') not in existing_any and quiz.get('id') in overdue_quiz_ids]
    if not missing:
        return quizzes

    insert_records = []
    for quiz in missing:
        out_of = quiz.get('questions') or 0
        insert_records.append({
            'quiz_id': quiz.get('id'),
            'student_id': student_id,
            'score': 0,
            'out_of': out_of,
            'grade': 'F',
            'passed': False,
            'status': 'missed',
            'attempted_at': datetime.now(timezone.utc).isoformat(),
        })

    if insert_records:
        insert_response = _safe_insert('quiz_attempts', insert_records)
        if supabase_failed(insert_response):
            raise HTTPException(status_code=502, detail=supabase_error_message(insert_response, 'Supabase insert overdue quiz attempts failed'))

    return quizzes


def _finalize_expired_in_progress_for_student(student_id: str) -> None:
    """Best-effort cleanup: whenever a student's dashboard is polled, close out any of
    their in-progress attempts whose quiz has passed its close_date, or whose personal
    time limit has elapsed, so quizzes/attempts don't stay open indefinitely just because
    nobody explicitly submitted or the tab was abandoned. Never raises — this is a
    housekeeping side-effect of read endpoints, not something that should break them.
    """
    if not student_id:
        return

    now = datetime.now(timezone.utc)
    _finalize_in_progress_attempts(student_id, now)
    _finalize_never_attempted_expired_quizzes(student_id, now)


def _finalize_in_progress_attempts(student_id: str, now: datetime) -> None:
    try:
        response = supabase.table('quiz_attempts').select('*').eq('student_id', student_id).eq('status', 'in_progress').execute()
    except _postgrest_exceptions.APIError as exc:
        if _is_missing_table_error(exc):
            return
        return
    except Exception:
        return

    if supabase_failed(response) or not response.data:
        return

    in_progress = response.data
    quiz_ids = [qid for qid in {a.get('quiz_id') for a in in_progress} if qid is not None]
    if not quiz_ids:
        return

    try:
        quizzes_resp = supabase.table('quizzes').select('*').in_('id', quiz_ids).execute()
    except Exception:
        return
    if supabase_failed(quizzes_resp):
        return
    quizzes_by_id = {q['id']: q for q in (quizzes_resp.data or []) if q.get('id') is not None}

    for attempt in in_progress:
        quiz = quizzes_by_id.get(attempt.get('quiz_id'))
        if not quiz:
            continue
        close_dt = _parse_datetime(quiz.get('close_date')) or _parse_datetime(quiz.get('due_date'))
        time_limit_minutes = int(quiz.get('time_limit') or 0)
        started_at = _parse_datetime(attempt.get('attempted_at'))

        past_close = bool(close_dt and now > close_dt)
        timed_out = bool(time_limit_minutes and started_at and (now - started_at).total_seconds() > time_limit_minutes * 60 + TIME_LIMIT_GRACE_SECONDS)

        if past_close or timed_out:
            try:
                supabase.table('quiz_attempts').update({
                    'status': 'missed',
                    'score': 0,
                    'grade': 'F',
                    'passed': False,
                    'attempted_at': now.isoformat(),
                }).eq('id', attempt['id']).execute()
            except Exception:
                continue


def _finalize_never_attempted_expired_quizzes(student_id: str, now: datetime) -> None:
    """A quiz the student was eligible for but never even opened shouldn't just vanish
    once it expires — it should land in their completed history as a missed (0-score)
    result, the same visible consequence as failing it, so it's honestly reflected in
    their record instead of silently disappearing.
    """
    try:
        student_level, student_program = _student_level_program(student_id)
    except Exception:
        return

    try:
        all_quizzes_resp = supabase.table('quizzes').select('*').execute()
    except Exception:
        return
    if supabase_failed(all_quizzes_resp):
        return

    try:
        allowed_codes = set(_fetch_allowed_course_codes(student_level, student_program)) if student_level else set()
    except Exception:
        return
    target_digit = _extract_level_digit(student_level) if student_level else None

    expired_eligible: dict[int, dict[str, Any]] = {}
    for quiz in (all_quizzes_resp.data or []):
        if str(quiz.get('status') or '').lower() == 'draft':
            continue
        if not _quiz_is_expired(quiz, now):
            continue
        qid = quiz.get('id')
        if qid is None:
            continue

        if student_level:
            course_code = str(quiz.get('course') or '')
            title = str(quiz.get('title') or '')
            clean_q = _clean_code_str(course_code)
            q_digit = _extract_level_digit(course_code) or _extract_level_digit(title)
            eligible = (
                course_code in allowed_codes
                or clean_q in allowed_codes
                or (target_digit is not None and q_digit is not None and target_digit == q_digit)
            )
            if not eligible:
                continue

        expired_eligible[qid] = quiz

    if not expired_eligible:
        return

    try:
        existing_resp = supabase.table('quiz_attempts').select('quiz_id').eq('student_id', student_id).in_('quiz_id', list(expired_eligible.keys())).execute()
    except Exception:
        return
    if supabase_failed(existing_resp):
        return
    existing_quiz_ids = {a.get('quiz_id') for a in (existing_resp.data or [])}

    missing_ids = [qid for qid in expired_eligible if qid not in existing_quiz_ids]
    if not missing_ids:
        return

    insert_records = [{
        'quiz_id': qid,
        'student_id': student_id,
        'score': 0,
        'out_of': expired_eligible[qid].get('questions') or 0,
        'grade': 'F',
        'passed': False,
        'status': 'missed',
        'attempted_at': now.isoformat(),
    } for qid in missing_ids]

    try:
        insert_resp = _safe_insert('quiz_attempts', insert_records)
        if supabase_failed(insert_resp):
            print(f"Failed to record missed (never-attempted) quizzes for student {student_id}: {supabase_error_message(insert_resp, '')}")
    except Exception as exc:
        print(f"Failed to record missed (never-attempted) quizzes for student {student_id}: {exc}")


def _safe_insert(table_name: str, record: dict | list[dict]) -> Any:
    """Insert a record or list of records into `table_name`, removing any keys reported
    as missing by PostgREST and retrying. Returns the supabase response object.
    """
    payload = record
    # normalize to list for uniform processing
    is_list = isinstance(payload, list)
    rows = payload if is_list else [payload]
    while True:
        try:
            resp = supabase.table(table_name).insert(rows if is_list else rows[0]).execute()
            return resp
        except _postgrest_exceptions.APIError as exc:
            msg = str(exc)
            if _is_missing_table_error(exc):
                raise HTTPException(status_code=502, detail=f"Database schema not ready: table '{table_name}' is missing")
            m = _re.search(r"Could not find the '([^']+)' column of '([^']+)'", msg)
            if not m:
                raise
            missing_col = m.group(1)
            # remove the key from all rows and retry
            changed = False
            for r in rows:
                if missing_col in r:
                    r.pop(missing_col, None)
                    changed = True
            if not changed:
                # nothing to remove, re-raise
                raise
        except Exception:
            raise


def _safe_update(table_name: str, payload: dict, in_key: str | None = None, in_values: list | None = None) -> Any:
    body = payload.copy()
    while True:
        try:
            q = supabase.table(table_name).update(body)
            if in_key and in_values is not None:
                q = q.in_(in_key, in_values)
            resp = q.execute()
            return resp
        except Exception as exc:
            msg = str(exc)
            m = _re.search(r"Could not find the '([^']+)' column of '([^']+)'", msg)
            if not m:
                raise
            missing_col = m.group(1)
            if missing_col in body:
                body.pop(missing_col, None)
                continue
            raise


class GenerateQuizRequest(BaseModel):
    course: str
    material_ids: list[int] = Field(default_factory=list)
    question_count: int = Field(default=10, ge=1, le=50)
    difficulty: str = 'Mixed'
    tier: str = 'Foundational'
    generate_all_tiers: bool = False
    question_types: list[str] = Field(default_factory=lambda: ['MCQ', 'True/False', 'Fill in the Blank'])
    time_limit: int = Field(default=45, ge=5)
    passing_score: int = Field(default=60, ge=0, le=100)
    attempts: int = Field(default=1, ge=1)
    due_in_days: int = Field(default=7, ge=1, le=30)
    open_date: str | None = None
    close_date: str | None = None


class PublishQuizRequest(BaseModel):
    title: str | None = None
    course: str | None = None
    questions: list[dict[str, Any]] = Field(default_factory=list)
    time_limit: int | None = None
    passing_score: int = Field(default=60, ge=0, le=100)
    attempts: int = Field(default=1, ge=1)
    due_date: str | None = None
    difficulty: str = 'Mixed'
    tier: str = 'Foundational'
    open_date: str | None = None
    close_date: str | None = None
    material_id: int | None = None
    material_ids: list[int] = Field(default_factory=list)
    quizzes: list[dict[str, Any]] | None = None


class QuizSubmissionRequest(BaseModel):
    student_id: str
    answers: dict[str, str]


def _get_student_attempts(quiz_id: int, student_id: str) -> list[dict[str, Any]]:
    try:
        response = supabase.table('quiz_attempts').select('*').eq('quiz_id', quiz_id).eq('student_id', student_id).execute()
    except _postgrest_exceptions.APIError as exc:
        if _is_missing_table_error(exc):
            return []
        raise
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase fetch student quiz attempts failed'))
    return response.data or []


def _fetch_student_attempt_counts(student_id: str) -> dict[int, int]:
    """Returns count of *completed or missed* attempts per quiz for the given student.
    In-progress attempts (score=0, not yet submitted) are intentionally excluded so
    that students who opened a quiz can still submit without being blocked.
    """
    try:
        response = supabase.table('quiz_attempts').select('quiz_id,status').eq('student_id', student_id).execute()
    except _postgrest_exceptions.APIError as exc:
        if _is_missing_table_error(exc):
            return {}
        raise
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase fetch student quiz attempts failed'))

    counts: dict[int, int] = {}
    for attempt in response.data or []:
        # Only count finalized attempts; ignore in_progress records
        if attempt.get('status') in ('completed', 'missed'):
            quiz_id = attempt.get('quiz_id')
            if quiz_id is not None:
                counts[quiz_id] = counts.get(quiz_id, 0) + 1
    return counts


def _extract_level_digit(text: str | None) -> str | None:
    if not text:
        return None
    m = _re.search(r'\d{3}', str(text))
    return m.group(0)[0] if m else None


def _clean_code_str(s: str | None) -> str:
    return _re.sub(r'[^a-zA-Z0-9]', '', s or '').lower()


def _fetch_allowed_course_codes(level: str | None, program: str | None) -> list[str]:
    if not level:
        return []
    response = supabase.table('courses').select('code', 'title', 'level', 'program').execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase fetch allowed course codes failed'))

    target_digit = _extract_level_digit(level)
    normalized_program = program.strip().lower() if program else ''
    allowed_codes: set[str] = set()

    for course in response.data or []:
        course_program = str(course.get('program') or '').strip().lower()
        if normalized_program and course_program and course_program != normalized_program:
            continue

        c_level = str(course.get('level') or '')
        c_code = str(course.get('code') or '')
        c_title = str(course.get('title') or '')
        c_digit = _extract_level_digit(c_level) or _extract_level_digit(c_code)

        if (level.lower() in c_level.lower()) or (c_level.lower() in level.lower()) or (target_digit and c_digit and target_digit == c_digit):
            if c_code:
                allowed_codes.add(c_code)
                allowed_codes.add(_clean_code_str(c_code))
            if c_title:
                allowed_codes.add(c_title)

    return list(allowed_codes)


def _student_level_program(student_id: str) -> tuple[str | None, str | None]:
    response = supabase.table('users').select('level,program,role').eq('id', student_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase fetch student profile failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Student not found')
    student = response.data[0]
    if student.get('role') != 'student':
        raise HTTPException(status_code=403, detail='User is not a student')
    return student.get('level'), student.get('program')


def _is_course_allowed_for_student(course_code: str, level: str | None, program: str | None, quiz_title: str = '') -> bool:
    if not level:
        return True
    allowed_codes = set(_fetch_allowed_course_codes(level, program))
    target_digit = _extract_level_digit(level)
    q_digit = _extract_level_digit(course_code) or _extract_level_digit(quiz_title)
    clean_q = _clean_code_str(course_code)
    return (
        course_code in allowed_codes
        or clean_q in allowed_codes
        or (target_digit is not None and q_digit is not None and target_digit == q_digit)
    )



def _normalize_question_types(question_types: list[str]) -> list[str]:
    valid_types = ['MCQ', 'True/False', 'Fill in the Blank', 'Short Answer']
    normalized = [qt for qt in question_types if qt in valid_types]
    return normalized or ['MCQ', 'True/False', 'Fill in the Blank']


@router.post('/generate')
def generate_quiz(payload: GenerateQuizRequest, _claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    ensure_supabase_enabled()
    try:
        if payload.material_ids:
            response = supabase.table('materials').select('*').in_('id', payload.material_ids).execute()
        else:
            response = supabase.table('materials').select('*').eq('course', payload.course).execute()
        if supabase_failed(response):
            raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase list materials failed'))
        materials = response.data or []
        if not materials:
            raise HTTPException(status_code=404, detail='No processed materials found for this course')

        # Inherit course from material if specified
        mat_course = next((m.get('course') for m in materials if m.get('course')), None)
        effective_course = mat_course or payload.course

        raw_titles = [material.get('name', '') for material in materials]
        material_titles = [_clean_material_title(title) for title in raw_titles]
        material_texts = []
        for material in materials:
            path = material.get('path', '') or material.get('file_path', '')
            text = ''
            if path and Path(path).exists():
                text = _extract_text_from_file(path)
            if not text and material.get('name'):
                mat_name = material.get('name', '')
                c_name = material.get('course') or effective_course
                local_dir = UPLOAD_DIR / c_name
                if local_dir.exists():
                    for f in local_dir.iterdir():
                        if f.is_file() and (f.name == mat_name or mat_name in f.name):
                            text = _extract_text_from_file(str(f))
                            break
            if not text:
                text = material.get('text_content', '') or material.get('name', '')
            material_texts.append(text)
        summarized_titles = [_summarize_material_text(text, raw_titles[i]) for i, text in enumerate(material_texts)]

        question_types = _normalize_question_types(payload.question_types)
        normalized_tier = _normalize_tier(payload.tier)
        now = datetime.now(timezone.utc)
        open_dt = _parse_datetime(payload.open_date) if payload.open_date else now
        close_dt = _parse_datetime(payload.close_date) if payload.close_date else now + timedelta(days=payload.due_in_days)
        if open_dt is None:
            open_dt = now
        if close_dt is None:
            close_dt = now + timedelta(days=payload.due_in_days)
        if close_dt <= open_dt:
            close_dt = open_dt + timedelta(days=payload.due_in_days)

        base_window = close_dt - open_dt
        if base_window.total_seconds() <= 0:
            base_window = timedelta(days=payload.due_in_days)

        def build_quiz_set(tier: str, offset: int) -> dict[str, Any]:
            effective_difficulty = payload.difficulty if payload.difficulty != 'Mixed' else _tier_default_difficulty(tier)
            questions: list[dict[str, Any]] = []
            for index in range(payload.question_count):
                question_type = question_types[index % len(question_types)]
                questions.append(_generate_question(effective_course, summarized_titles, effective_difficulty, question_type, index + offset * 1000))
            tier_open = open_dt + offset * base_window
            tier_close = tier_open + base_window
            return {
                'title': f"AI Generated {tier} Quiz for {effective_course}",
                'course': effective_course,
                'questions': len(questions),
                'time_limit': payload.time_limit,
                'passing_score': payload.passing_score,
                'attempts': payload.attempts,
                'due_date': _format_datetime(tier_close),
                'open_date': _format_datetime(tier_open),
                'close_date': _format_datetime(tier_close),
                'difficulty': effective_difficulty,
                'tier': tier,
                'material_ids': [material.get('id') for material in materials],
                'questions': questions,
            }

        if payload.generate_all_tiers:
            quiz_sets = []
            questions_by_tier: dict[str, list[dict[str, Any]]] = {}
            draft_quiz_ids: dict[str, int] = {}
            for offset, tier in enumerate(['Foundational', 'Intermediate', 'Mastery']):
                quiz_set = build_quiz_set(tier, offset)
                tier_questions = quiz_set.pop('questions')

                # Persist as a draft immediately, rather than only on publish, so the bank
                # shows up in the lecturer's Question Banks archive right away — previously
                # a generated-but-not-yet-published bank existed only in React state and
                # vanished (with no trace anywhere) the moment the lecturer navigated away.
                draft_id: int | None = None
                try:
                    quiz_record = {
                        'title': quiz_set['title'],
                        'course': quiz_set['course'],
                        'questions': len(tier_questions),
                        'time_limit': quiz_set['time_limit'],
                        'passing_score': quiz_set['passing_score'],
                        'attempts': quiz_set['attempts'],
                        'due_date': quiz_set['due_date'],
                        'open_date': None,
                        'close_date': None,
                        'status': 'draft',
                        'difficulty': quiz_set['difficulty'],
                        'tier': tier,
                        'material_ids': quiz_set['material_ids'],
                    }
                    quiz_resp = _safe_insert('quizzes', quiz_record)
                    if not supabase_failed(quiz_resp) and quiz_resp.data:
                        draft_id = quiz_resp.data[0]['id']
                        question_rows = [{
                            'quiz_id': draft_id,
                            'question': q.get('question', ''),
                            'options': q.get('options', []),
                            'correct': q.get('answer', ''),
                        } for q in tier_questions]
                        q_resp = _safe_insert('quiz_questions', question_rows)
                        if not supabase_failed(q_resp) and q_resp.data:
                            for question, row in zip(tier_questions, q_resp.data):
                                question['id'] = row['id']
                except Exception as persist_err:
                    # A persistence hiccup here shouldn't block the lecturer from reviewing
                    # and publishing straight from the in-memory generation result.
                    print(f"Draft quiz persistence failed for {tier}: {persist_err}")

                if draft_id is not None:
                    draft_quiz_ids[tier] = draft_id
                quiz_set['id'] = draft_id
                questions_by_tier[tier] = tier_questions
                quiz_sets.append(quiz_set)

            return {
                'quiz_sets': quiz_sets,
                'questions_by_tier': questions_by_tier,
                'materials': materials,
                'draft_quiz_ids': draft_quiz_ids,
            }

        questions: list[dict[str, Any]] = []
        for index in range(payload.question_count):
            question_type = question_types[index % len(question_types)]
            questions.append(_generate_question(effective_course, summarized_titles, payload.difficulty if payload.difficulty != 'Mixed' else _tier_default_difficulty(normalized_tier), question_type, index))

        return {
            'quiz': {
                'title': f"AI Generated {normalized_tier} Quiz for {effective_course}",
                'course': effective_course,
                'questions': len(questions),
                'time_limit': payload.time_limit,
                'passing_score': payload.passing_score,
                'attempts': payload.attempts,
                'due_date': _format_datetime(close_dt),
                'open_date': _format_datetime(open_dt),
                'close_date': _format_datetime(close_dt),
                'difficulty': payload.difficulty if payload.difficulty != 'Mixed' else _tier_default_difficulty(normalized_tier),
                'tier': normalized_tier,
                'material_ids': [material.get('id') for material in materials],
            },
            'questions': questions,
            'materials': materials,
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"Quiz generation exception: {e}")
        raise HTTPException(status_code=500, detail=f"Quiz generation error: {str(e)}")



@router.post('/publish')
def publish_quiz(payload: PublishQuizRequest, _claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    ensure_supabase_enabled()

    # If multiple quiz sets are provided (tiered generation), publish each separately
    if payload.quizzes:
        published: list[dict[str, Any]] = []
        all_questions: dict[int, list[dict[str, Any]]] = {}
        for q in payload.quizzes:
            due_dt = _parse_datetime(q.get('due_date'))
            if due_dt is None:
                raise HTTPException(status_code=400, detail='Quiz due date must be a valid ISO datetime string')

            # SECURITY FIX: if no open_date provided, quiz stays locked ('scheduled')
            open_dt = _parse_datetime(q.get('open_date')) if q.get('open_date') else None
            close_dt = _parse_datetime(q.get('close_date')) if q.get('close_date') else due_dt
            if close_dt is None:
                raise HTTPException(status_code=400, detail='Quiz close date must be a valid ISO datetime string')
            if open_dt and close_dt < open_dt:
                raise HTTPException(status_code=400, detail='Quiz close date must be after open date')

            now = datetime.now(timezone.utc)
            if open_dt is None:
                status = 'scheduled'  # No open_date means locked until lecturer sets one
            elif open_dt <= now <= close_dt:
                status = 'available'
            elif open_dt > now:
                status = 'scheduled'
            else:
                status = 'closed'

            normalized_tier = _normalize_tier(q.get('tier') or payload.tier)
            derived_course = q.get('course') or payload.course
            derived_title = q.get('title') or payload.title or f"AI Generated {normalized_tier} Quiz for {derived_course}"

            quiz_record = {
                'title': derived_title,
                'course': q.get('course') or payload.course,
                'questions': len(q.get('questions', []) or []),
                'time_limit': q.get('time_limit') if q.get('time_limit') is not None else payload.time_limit,
                'passing_score': q.get('passing_score', payload.passing_score),
                'attempts': q.get('attempts', payload.attempts),
                'due_date': q.get('due_date'),
                'open_date': _format_datetime(open_dt) if open_dt else None,
                'close_date': _format_datetime(close_dt),
                'status': status,
                'difficulty': q.get('difficulty', payload.difficulty),
                'tier': normalized_tier,
                'material_id': q.get('material_id') or payload.material_id,
                'material_ids': q.get('material_ids', payload.material_ids or []),
            }

            draft_id = q.get('id')
            if draft_id is not None:
                # This bank was already persisted as a draft at generation time — publishing
                # finalizes the SAME row (schedule + status) rather than inserting a duplicate.
                quiz_response = _safe_update('quizzes', quiz_record, 'id', [draft_id])
                if supabase_failed(quiz_response):
                    raise HTTPException(status_code=502, detail=supabase_error_message(quiz_response, 'Supabase update draft quiz failed'))
                if not quiz_response.data:
                    raise HTTPException(status_code=502, detail='Draft quiz not found')
                created_quiz = quiz_response.data[0]

                # Replace the draft's questions with the lecturer's final (post-review,
                # post-approval) set rather than appending duplicates.
                delete_resp = supabase.table('quiz_questions').delete().eq('quiz_id', draft_id).execute()
                if supabase_failed(delete_resp):
                    raise HTTPException(status_code=502, detail=supabase_error_message(delete_resp, 'Supabase clear draft quiz questions failed'))
            else:
                quiz_response = _safe_insert('quizzes', quiz_record)
                if supabase_failed(quiz_response):
                    raise HTTPException(status_code=502, detail=supabase_error_message(quiz_response, 'Supabase insert quiz failed'))
                if not quiz_response.data:
                    raise HTTPException(status_code=502, detail='Failed to create quiz')
                created_quiz = quiz_response.data[0]

            question_rows = []
            for question in q.get('questions', []):
                question_rows.append({
                    'quiz_id': created_quiz['id'],
                    'question': question.get('question', ''),
                    'options': question.get('options', []),
                    'correct': question.get('answer', ''),
                })
            questions_response = _safe_insert('quiz_questions', question_rows)
            if supabase_failed(questions_response):
                raise HTTPException(status_code=502, detail=supabase_error_message(questions_response, 'Supabase insert quiz questions failed'))

            m_ids = q.get('material_ids') or payload.material_ids or []
            if m_ids:
                update_response = _safe_update('materials', {'quiz_generated': True}, 'id', m_ids)
                if supabase_failed(update_response):
                    raise HTTPException(status_code=502, detail=supabase_error_message(update_response, 'Supabase update materials failed'))

            published.append(created_quiz)
            all_questions[created_quiz['id']] = questions_response.data or []

        return {'quizzes': published, 'questions': all_questions}

    # Single-quiz publish (backwards compatible)
    raw_due = payload.due_date or payload.close_date or (datetime.now(timezone.utc) + timedelta(days=7)).isoformat()
    due_dt = _parse_datetime(raw_due) or (datetime.now(timezone.utc) + timedelta(days=7))

    # SECURITY FIX: if no open_date is provided, keep quiz as 'scheduled' (locked)
    # Never default open_date to now() — lecturers must explicitly set an opening time
    open_dt = _parse_datetime(payload.open_date) if payload.open_date else None
    close_dt = _parse_datetime(payload.close_date) if payload.close_date else due_dt
    if close_dt is None:
        raise HTTPException(status_code=400, detail='Quiz close date must be a valid ISO datetime string')
    if open_dt is None:
        open_dt = datetime.now(timezone.utc)
    if close_dt < open_dt:
        raise HTTPException(status_code=400, detail='Quiz close date must be after open date')

    now = datetime.now(timezone.utc)
    if open_dt <= now <= close_dt:
        status = 'available'
    elif open_dt > now:
        status = 'scheduled'
    else:
        status = 'closed'
    normalized_tier = _normalize_tier(payload.tier)

    # Derive title when not provided to reflect the publishing course
    derived_course = payload.course
    derived_title = payload.title or f"AI Generated {normalized_tier} Quiz for {derived_course}"

    quiz_record = {
        'title': derived_title,
        'course': payload.course,
        'questions': len(payload.questions),
        'time_limit': payload.time_limit,
        'passing_score': payload.passing_score,
        'attempts': payload.attempts,
        'due_date': payload.due_date,
        'open_date': _format_datetime(open_dt) if open_dt else None,
        'close_date': _format_datetime(close_dt),
        'status': status,
        'difficulty': payload.difficulty,
        'tier': normalized_tier,
        'material_id': payload.material_id,
        'material_ids': payload.material_ids,
    }
    quiz_response = _safe_insert('quizzes', quiz_record)
    if supabase_failed(quiz_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(quiz_response, 'Supabase insert quiz failed'))
    if not quiz_response.data:
        raise HTTPException(status_code=502, detail='Failed to create quiz')

    quiz = quiz_response.data[0]
    question_rows = []
    for question in payload.questions:
        question_rows.append({
            'quiz_id': quiz['id'],
            'question': question.get('question', ''),
            'options': question.get('options', []),
            'correct': question.get('answer', ''),
        })
    questions_response = _safe_insert('quiz_questions', question_rows)
    if supabase_failed(questions_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(questions_response, 'Supabase insert quiz questions failed'))

    if payload.material_ids:
        update_response = _safe_update('materials', {'quiz_generated': True}, 'id', payload.material_ids)
        if supabase_failed(update_response):
            raise HTTPException(status_code=502, detail=supabase_error_message(update_response, 'Supabase update materials failed'))

    return {'quiz': quiz, 'questions': questions_response.data or []}



@router.delete('/all')
def delete_all_quizzes(_claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    ensure_supabase_enabled()

    # Delete all quiz attempts first (foreign key dependency)
    try:
        attempts_resp = supabase.table('quiz_attempts').delete().neq('id', 0).execute()
        if supabase_failed(attempts_resp):
            raise HTTPException(status_code=502, detail=supabase_error_message(attempts_resp, 'Failed to delete quiz attempts'))
    except _postgrest_exceptions.APIError as exc:
        if not _is_missing_table_error(exc):
            raise HTTPException(status_code=502, detail=str(exc))

    # Delete all quiz questions
    try:
        questions_resp = supabase.table('quiz_questions').delete().neq('id', 0).execute()
        if supabase_failed(questions_resp):
            raise HTTPException(status_code=502, detail=supabase_error_message(questions_resp, 'Failed to delete quiz questions'))
    except _postgrest_exceptions.APIError as exc:
        if not _is_missing_table_error(exc):
            raise HTTPException(status_code=502, detail=str(exc))

    # Delete all quizzes
    quizzes_resp = supabase.table('quizzes').delete().neq('id', 0).execute()
    if supabase_failed(quizzes_resp):
        raise HTTPException(status_code=502, detail=supabase_error_message(quizzes_resp, 'Failed to delete quizzes'))

    return {'message': 'All quizzes, questions, and attempts have been deleted successfully.'}


class QuizStartRequest(BaseModel):
    student_id: str


@router.post('/{quiz_id}/start')
def start_quiz(quiz_id: int, payload: QuizStartRequest, claims: dict = Depends(get_current_claims)) -> dict[str, Any]:
    """Called the moment a student opens the quiz. Records an in-progress attempt
    with score=0 immediately, so that even if the student closes without submitting
    a score of 0 is persisted."""
    _verify_acting_as_self(claims, payload.student_id)
    ensure_supabase_enabled()
    response = supabase.table('quizzes').select('*').eq('id', quiz_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase get quiz failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Quiz not found')
    quiz = response.data[0]

    student_level, student_program = _student_level_program(payload.student_id)
    if not _is_course_allowed_for_student(quiz.get('course', ''), student_level, student_program):
        raise HTTPException(status_code=403, detail='Student is not allowed to access this quiz')

    now = datetime.now(timezone.utc)
    open_dt = _parse_datetime(quiz.get('open_date'))
    close_dt = _parse_datetime(quiz.get('close_date'))
    if open_dt and now < open_dt:
        raise HTTPException(status_code=403, detail='Quiz is not open yet')
    if close_dt and now > close_dt:
        raise HTTPException(status_code=403, detail='Quiz deadline has passed')

    attempts = _get_student_attempts(quiz_id, payload.student_id)

    # Already has a completed attempt — block
    completed = [a for a in attempts if a.get('status') == 'completed']
    if completed:
        raise HTTPException(status_code=403, detail='No attempts remaining for this quiz')

    time_limit_minutes = int(quiz.get('time_limit') or 0)

    # Already has an in-progress attempt — return it (idempotent), unless its
    # personal time limit has already elapsed, in which case it is finalized as missed.
    in_progress = [a for a in attempts if a.get('status') == 'in_progress']
    if in_progress:
        existing = in_progress[0]
        started_at = _parse_datetime(existing.get('attempted_at'))
        if time_limit_minutes and started_at and (now - started_at).total_seconds() > time_limit_minutes * 60 + TIME_LIMIT_GRACE_SECONDS:
            supabase.table('quiz_attempts').update({
                'status': 'missed',
                'score': 0,
                'grade': 'F',
                'passed': False,
                'attempted_at': now.isoformat(),
            }).eq('id', existing['id']).execute()
            raise HTTPException(status_code=403, detail='Time limit exceeded for your previous attempt; it has been recorded as missed.')

        expires_at = _format_datetime(started_at + timedelta(minutes=time_limit_minutes)) if (time_limit_minutes and started_at) else None
        return {
            'attempt': existing,
            'already_started': True,
            'time_limit': time_limit_minutes,
            'started_at': existing.get('attempted_at'),
            'expires_at': expires_at,
        }

    # Fetch total question count
    qcount_resp = supabase.table('quiz_questions').select('id').eq('quiz_id', quiz_id).execute()
    total_questions = len(qcount_resp.data or []) if not supabase_failed(qcount_resp) else 0

    # Record in-progress attempt with score 0 immediately
    started_at_iso = now.isoformat()
    attempt_record = {
        'quiz_id': quiz_id,
        'student_id': payload.student_id,
        'score': 0,
        'out_of': total_questions,
        'grade': 'F',
        'passed': False,
        'status': 'in_progress',
        'attempted_at': started_at_iso,
    }
    insert_response = _safe_insert('quiz_attempts', attempt_record)
    if supabase_failed(insert_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(insert_response, 'Supabase insert quiz attempt failed'))

    expires_at = _format_datetime(now + timedelta(minutes=time_limit_minutes)) if time_limit_minutes else None
    return {
        'attempt': insert_response.data[0] if insert_response.data else attempt_record,
        'already_started': False,
        'time_limit': time_limit_minutes,
        'started_at': started_at_iso,
        'expires_at': expires_at,
    }


@router.post('/{quiz_id}/submit')
def submit_quiz(quiz_id: int, payload: QuizSubmissionRequest, claims: dict = Depends(get_current_claims)) -> dict[str, Any]:
    _verify_acting_as_self(claims, payload.student_id)
    ensure_supabase_enabled()
    response = supabase.table('quizzes').select('*').eq('id', quiz_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase get quiz failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Quiz not found')
    quiz = response.data[0]

    student_level, student_program = _student_level_program(payload.student_id)
    if not _is_course_allowed_for_student(quiz.get('course', ''), student_level, student_program):
        raise HTTPException(status_code=403, detail='Student is not allowed to access this quiz')

    now = datetime.now(timezone.utc)
    open_dt = _parse_datetime(quiz.get('open_date'))
    close_dt = _parse_datetime(quiz.get('close_date'))
    if open_dt and now < open_dt:
        raise HTTPException(status_code=403, detail='Quiz is not open yet')
    if close_dt and now > close_dt:
        _finalize_overdue_attempts([quiz], payload.student_id)
        raise HTTPException(status_code=403, detail='Quiz deadline has passed')

    attempts = _get_student_attempts(quiz_id, payload.student_id)

    # Block if already has a completed attempt
    completed_attempts = [a for a in attempts if a.get('status') == 'completed']
    if completed_attempts:
        raise HTTPException(status_code=403, detail='No attempts remaining for this quiz')

    # Find the in-progress attempt to update (started via /start)
    in_progress_attempts = [a for a in attempts if a.get('status') == 'in_progress']
    in_progress_attempt_id = in_progress_attempts[0].get('id') if in_progress_attempts else None

    # If somehow they submit without starting (legacy), block if any attempt exists
    if not in_progress_attempt_id and len(attempts) >= 1:
        raise HTTPException(status_code=403, detail='No attempts remaining for this quiz')

    # SECURITY: enforce the per-attempt time limit server-side, not just via the
    # client-side countdown, so a tampered/stalled client cannot buy extra time.
    time_limit_minutes = int(quiz.get('time_limit') or 0)
    if in_progress_attempt_id and time_limit_minutes:
        started_at = _parse_datetime(in_progress_attempts[0].get('attempted_at'))
        if started_at and (now - started_at).total_seconds() > time_limit_minutes * 60 + TIME_LIMIT_GRACE_SECONDS:
            supabase.table('quiz_attempts').update({
                'status': 'missed',
                'score': 0,
                'grade': 'F',
                'passed': False,
                'attempted_at': now.isoformat(),
            }).eq('id', in_progress_attempt_id).execute()
            raise HTTPException(status_code=403, detail='Time limit exceeded. Your attempt has been recorded as missed with a score of 0.')

    questions_response = supabase.table('quiz_questions').select('*').eq('quiz_id', quiz_id).execute()
    if supabase_failed(questions_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(questions_response, 'Supabase get quiz questions failed'))
    raw_questions = questions_response.data or []

    import hashlib
    seed_key = f"{quiz_id}_{payload.student_id or 'random'}"
    stable_seed = int(hashlib.md5(seed_key.encode()).hexdigest(), 16)
    rng = random.Random(stable_seed)

    questions = []
    for q in raw_questions:
        q_copy = dict(q)
        opts = q_copy.get('options')
        if isinstance(opts, list) and len(opts) > 1:
            opts_copy = list(opts)
            rng.shuffle(opts_copy)
            q_copy['options'] = opts_copy
        questions.append(q_copy)

    rng.shuffle(questions)
    correct = 0

    def _answers_match(expected: str, given: str, q_type: str) -> bool:
        """Return True if the student's answer is correct.
        - Always case-insensitive.
        - For MCQ / True-False: exact match only.
        - For Fill in the Blank / Short Answer: also accept partial / contained matches
          so 'Halt' matches 'halt' and semantically similar phrases pass.
        """
        exp = expected.strip().lower()
        ans = given.strip().lower()
        if not ans:
            return False
        if exp == ans:
            return True
        # For text-entry types allow partial match
        text_types = {'fill in the blank', 'fill-in-the-blank', 'short answer', 'short_answer', 'fitb'}
        if str(q_type).strip().lower() in text_types:
            # Pass if one is a substring of the other
            if exp in ans or ans in exp:
                return True
            # Pass if student answer contains every word of the expected answer
            exp_words = set(exp.split())
            ans_words = set(ans.split())
            if exp_words and exp_words.issubset(ans_words):
                return True
        return False

    for idx, question in enumerate(questions):
        answer = payload.answers.get(str(idx), '')
        q_type = question.get('type', 'MCQ')
        if _answers_match(question.get('correct', ''), answer, q_type):
            correct += 1


    total_questions = len(questions)
    score = round((correct / total_questions) * 100) if total_questions else 0
    passed = score >= (quiz.get('passing_score') or 0)
    grade = _grade_from_score(score)

    attempt_record = {
        'quiz_id': quiz_id,
        'student_id': payload.student_id,
        'score': score,
        'out_of': total_questions,
        'grade': grade,
        'passed': passed,
        'status': 'completed',
        'attempted_at': datetime.now(timezone.utc).isoformat(),
    }

    if in_progress_attempt_id:
        # Update the existing in-progress record to completed
        update_resp = supabase.table('quiz_attempts').update({
            'score': score,
            'out_of': total_questions,
            'grade': grade,
            'passed': passed,
            'status': 'completed',
            'attempted_at': datetime.now(timezone.utc).isoformat(),
        }).eq('id', in_progress_attempt_id).execute()
        if supabase_failed(update_resp):
            raise HTTPException(status_code=502, detail=supabase_error_message(update_resp, 'Supabase update quiz attempt failed'))
    else:
        # Legacy path: no in-progress record, insert directly
        insert_response = _safe_insert('quiz_attempts', attempt_record)
        if supabase_failed(insert_response):
            raise HTTPException(status_code=502, detail=supabase_error_message(insert_response, 'Supabase insert quiz attempt failed'))

    return {'attempt': attempt_record, 'quiz': quiz}


class QuizForfeitRequest(BaseModel):
    student_id: str


@router.post('/{quiz_id}/forfeit')
def forfeit_quiz(quiz_id: int, payload: QuizForfeitRequest, claims: dict | None = Depends(get_current_claims_optional)) -> dict[str, Any]:
    """Called the instant a student navigates away from / backgrounds an in-progress
    attempt (tab switch, app switch, closing the window). Leaving the assessment screen
    is how a student would copy questions out to a third-party solver, so the attempt is
    finalized immediately as a 0-score miss rather than left open for them to resume.

    Identity is only checked when a token IS present: this is fired via
    `navigator.sendBeacon` on the frontend, which cannot attach an Authorization header by
    browser spec, so requiring auth here would silently break the anti-cheat forfeit itself.
    """
    if claims is not None:
        _verify_acting_as_self(claims, payload.student_id)
    ensure_supabase_enabled()
    if not payload.student_id:
        raise HTTPException(status_code=400, detail='student_id is required')

    attempts = _get_student_attempts(quiz_id, payload.student_id)
    in_progress = [a for a in attempts if a.get('status') == 'in_progress']
    if not in_progress:
        # Already completed/missed, or never started — nothing to forfeit.
        return {'status': 'no_active_attempt'}

    attempt = in_progress[0]
    update_resp = supabase.table('quiz_attempts').update({
        'status': 'missed',
        'score': 0,
        'grade': 'F',
        'passed': False,
        'attempted_at': datetime.now(timezone.utc).isoformat(),
    }).eq('id', attempt['id']).execute()
    if supabase_failed(update_resp):
        raise HTTPException(status_code=502, detail=supabase_error_message(update_resp, 'Supabase forfeit quiz attempt failed'))

    return {'status': 'forfeited', 'quiz_id': quiz_id}


def _quiz_is_expired(quiz: dict[str, Any], now: datetime) -> bool:
    close_dt = _parse_datetime(quiz.get('close_date')) or _parse_datetime(quiz.get('due_date'))
    if close_dt and now > close_dt:
        return True
    return False


def _quiz_is_locked(quiz: dict[str, Any], now: datetime) -> bool:
    open_dt = _parse_datetime(quiz.get('open_date'))
    if open_dt and now < open_dt:
        return True
    # Also lock if quiz has no open_date AND was explicitly scheduled/still a draft
    if not open_dt and str(quiz.get('status') or '').lower() in ('scheduled', 'draft'):
        return True
    return False


@router.get('/stats')
def get_quiz_stats() -> dict[str, Any]:
    """Institution-wide quiz completion snapshot for the admin dashboard. `/completed`
    is per-student (requires student_id) so it always returned empty for admin callers —
    this is the aggregate equivalent: how many of all published quizzes have received at
    least one completed attempt from any student."""
    ensure_supabase_enabled()
    quizzes_resp = supabase.table('quizzes').select('id').execute()
    if supabase_failed(quizzes_resp):
        raise HTTPException(status_code=502, detail=supabase_error_message(quizzes_resp, 'Supabase list quizzes failed'))
    total_quizzes = len(quizzes_resp.data or [])

    quizzes_with_completions = 0
    try:
        attempts_resp = supabase.table('quiz_attempts').select('quiz_id,status').eq('status', 'completed').execute()
        if not supabase_failed(attempts_resp):
            quizzes_with_completions = len({a['quiz_id'] for a in (attempts_resp.data or []) if a.get('quiz_id') is not None})
    except Exception:
        pass

    return {'total_quizzes': total_quizzes, 'quizzes_with_completions': quizzes_with_completions}


@router.get('')
def list_all_quizzes(course: str | None = None, lecturer: str | None = None, _claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    """Lecturer-facing management listing: every quiz for a course (or all of a lecturer's
    courses), with a live-computed status so a stale DB 'status' column never misleads a
    lecturer trying to figure out why a quiz is still locked/closed."""
    ensure_supabase_enabled()
    query = supabase.table('quizzes').select('*')
    if course:
        query = query.eq('course', course)
    elif lecturer:
        lecturer = lecturer.strip()
        courses_resp = supabase.table('courses').select('code').ilike('lecturer', lecturer).execute()
        if supabase_failed(courses_resp):
            raise HTTPException(status_code=502, detail=supabase_error_message(courses_resp, 'Supabase list lecturer courses failed'))
        course_codes = [c['code'] for c in (courses_resp.data or []) if c.get('code')]
        if not course_codes:
            return {'quizzes': []}
        query = query.in_('course', course_codes)
    response = query.execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase list quizzes failed'))

    now = datetime.now(timezone.utc)
    quizzes = []
    for quiz in response.data or []:
        q = dict(quiz)
        locked = _quiz_is_locked(q, now)
        expired = _quiz_is_expired(q, now)
        q['is_locked'] = locked
        q['is_closed'] = expired and not locked
        q['live_status'] = 'scheduled' if locked else ('closed' if expired else 'available')
        quizzes.append(q)
    quizzes.sort(key=lambda q: (q.get('course') or '', q.get('tier') or '', q.get('id') or 0))
    return {'quizzes': quizzes}


@router.get('/{quiz_id}/full')
def get_quiz_full(quiz_id: int, _claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    """Lecturer-only view of a quiz's complete question bank (including correct answers),
    with no lock/open-date restriction — used by the Question Banks archive to preview and
    export a hardcopy of what was generated, independent of whether students can see it yet."""
    ensure_supabase_enabled()
    response = supabase.table('quizzes').select('*').eq('id', quiz_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase get quiz failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Quiz not found')
    quiz = response.data[0]

    questions_response = supabase.table('quiz_questions').select('*').eq('quiz_id', quiz_id).execute()
    if supabase_failed(questions_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(questions_response, 'Supabase get quiz questions failed'))

    return {'quiz': quiz, 'questions': questions_response.data or []}


class UpdateQuizScheduleRequest(BaseModel):
    open_date: str | None = None
    close_date: str | None = None
    time_limit: int | None = Field(default=None, ge=5, le=180)
    passing_score: int | None = Field(default=None, ge=0, le=100)
    attempts: int | None = Field(default=None, ge=1)


@router.patch('/{quiz_id}/schedule')
def update_quiz_schedule(quiz_id: int, payload: UpdateQuizScheduleRequest, _claims: dict = Depends(require_roles('lecturer', 'admin', 'administrator'))) -> dict[str, Any]:
    """Lets a lecturer fix a quiz's open/close window (or duration/pass score/attempts)
    after it has already been published — e.g. to recover a quiz that got stuck 'scheduled'
    because of a bad date, without having to delete and regenerate the whole question bank."""
    ensure_supabase_enabled()
    response = supabase.table('quizzes').select('*').eq('id', quiz_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase get quiz failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Quiz not found')
    quiz = response.data[0]

    open_dt = _parse_datetime(payload.open_date) if payload.open_date is not None else _parse_datetime(quiz.get('open_date'))
    if payload.open_date is not None and open_dt is None:
        raise HTTPException(status_code=400, detail='Invalid open date')

    close_dt = _parse_datetime(payload.close_date) if payload.close_date is not None else _parse_datetime(quiz.get('close_date'))
    if payload.close_date is not None and close_dt is None:
        raise HTTPException(status_code=400, detail='Invalid close date')

    if open_dt and close_dt and close_dt <= open_dt:
        raise HTTPException(status_code=400, detail='Close date must be after open date')

    now = datetime.now(timezone.utc)
    if open_dt is None:
        status = 'scheduled'
    elif now < open_dt:
        status = 'scheduled'
    elif close_dt and now > close_dt:
        status = 'closed'
    else:
        status = 'available'

    update_payload: dict[str, Any] = {'status': status}
    if payload.open_date is not None:
        update_payload['open_date'] = _format_datetime(open_dt) if open_dt else None
    if payload.close_date is not None:
        update_payload['close_date'] = _format_datetime(close_dt) if close_dt else None
        if close_dt:
            update_payload['due_date'] = _format_datetime(close_dt)
    if payload.time_limit is not None:
        update_payload['time_limit'] = payload.time_limit
    if payload.passing_score is not None:
        update_payload['passing_score'] = payload.passing_score
    if payload.attempts is not None:
        update_payload['attempts'] = payload.attempts

    update_resp = _safe_update('quizzes', update_payload, 'id', [quiz_id])
    if supabase_failed(update_resp):
        raise HTTPException(status_code=502, detail=supabase_error_message(update_resp, 'Supabase update quiz schedule failed'))

    updated_quiz = (update_resp.data or [None])[0] or {**quiz, **update_payload}
    return {'quiz': updated_quiz}


@router.get('/available')
def list_available_quizzes(level: str | None = None, program: str | None = None, student_id: str | None = None) -> dict[str, Any]:
    ensure_supabase_enabled()
    if student_id:
        _finalize_expired_in_progress_for_student(student_id)
    response = supabase.table('quizzes').select('*').execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase list available quizzes failed'))

    all_quizzes = response.data or []
    target_digit = _extract_level_digit(level) if level else None
    allowed_codes = set(_fetch_allowed_course_codes(level, program)) if level else set()

    attempts_by_quiz: dict[int, int] = {}
    if student_id:
        attempts_by_quiz = _fetch_student_attempt_counts(student_id)

    now = datetime.now(timezone.utc)
    quizzes: list[dict[str, Any]] = []
    for quiz in all_quizzes:
        if _quiz_is_expired(quiz, now):
            continue
        # Drafts are persisted the moment a lecturer generates a bank (so it shows up in
        # their Question Banks archive immediately), but haven't been through the schedule
        # step yet — students must never see them, not even as a locked/upcoming entry.
        if str(quiz.get('status') or '').lower() == 'draft':
            continue

        quiz_dict = dict(quiz)
        is_locked = _quiz_is_locked(quiz_dict, now)
        if is_locked:
            quiz_dict['status'] = 'scheduled'
            quiz_dict['is_locked'] = True
            # Preserve the configured open_date so students can see it (even if in past/null)
            # open_date is already in quiz_dict from the DB
        else:
            quiz_dict['status'] = 'available'
            quiz_dict['is_locked'] = False

        if student_id:
            quiz_id = quiz_dict.get('id')
            if quiz_id is not None and attempts_by_quiz.get(quiz_id, 0) >= 1:
                continue

        if level:
            quiz_course = str(quiz_dict.get('course') or '')
            quiz_title = str(quiz_dict.get('title') or '')
            clean_q = _clean_code_str(quiz_course)
            q_digit = _extract_level_digit(quiz_course) or _extract_level_digit(quiz_title)

            matches_level = (
                quiz_course in allowed_codes
                or clean_q in allowed_codes
                or (target_digit is not None and q_digit is not None and target_digit == q_digit)
            )
            if not matches_level:
                continue

        quizzes.append(quiz_dict)

    # Security: strip actual question content from locked entries so students can't
    # retrieve questions/answers for a quiz that hasn't opened yet. Schedule/format
    # metadata (question count, time limit, pass score, due date) is safe to show —
    # it comes from the `quizzes` row, not the `quiz_questions` table, so exposing it
    # can't leak any question text or answers, and it's what students need to prepare.
    safe_quizzes = []
    for q in quizzes:
        if q.get('is_locked'):
            safe_q = {
                'id':             q.get('id'),
                'title':          q.get('title'),
                'course':         q.get('course'),
                'tier':           q.get('tier'),
                'status':         q.get('status'),
                'is_locked':      True,
                'open_date':      q.get('open_date'),
                'close_date':     q.get('close_date'),
                'due_date':       q.get('due_date'),
                'questions':      q.get('questions'),
                'time_limit':     q.get('time_limit'),
                'passing_score':  q.get('passing_score'),
                'attempts':       q.get('attempts'),
                # Intentionally omitted: actual question text, options, and answers.
            }
            safe_quizzes.append(safe_q)
        else:
            safe_quizzes.append(q)

    return {'quizzes': safe_quizzes}


@router.get('/completed')
def list_completed_quizzes(student_id: str | None = None, level: str | None = None, program: str | None = None) -> dict[str, Any]:
    ensure_supabase_enabled()
    if not student_id:
        return {'quizzes': []}

    _finalize_expired_in_progress_for_student(student_id)

    try:
        attempts_response = supabase.table('quiz_attempts').select('*').eq('student_id', student_id).execute()
    except Exception:
        return {'quizzes': []}
    if supabase_failed(attempts_response):
        return {'quizzes': []}
    attempts = attempts_response.data or []
    if not attempts:
        return {'quizzes': []}

    attempts_by_quiz: dict[int, dict[str, Any]] = {}
    for attempt in attempts:
        # Only include finalized attempts in the history view
        if attempt.get('status') == 'in_progress':
            continue
        quiz_id = attempt.get('quiz_id')
        if quiz_id is None:
            continue
        existing = attempts_by_quiz.get(quiz_id)
        if existing is None or (attempt.get('attempted_at') or '') > (existing.get('attempted_at') or ''):
            attempts_by_quiz[quiz_id] = attempt

    quiz_ids = list(attempts_by_quiz.keys())
    if not quiz_ids:
        return {'quizzes': []}

    all_quizzes = []
    chunk_size = 50
    for i in range(0, len(quiz_ids), chunk_size):
        chunk = quiz_ids[i:i + chunk_size]
        try:
            chunk_response = supabase.table('quizzes').select('*').in_('id', chunk).execute()
            if not supabase_failed(chunk_response) and chunk_response.data:
                all_quizzes.extend(chunk_response.data)
        except Exception as e:
            print(f"Error fetching quizzes chunk: {e}")

    quizzes_by_id = {quiz['id']: quiz for quiz in all_quizzes if quiz.get('id') is not None}
    allowed_codes = None
    if level:
        allowed_codes = set(_fetch_allowed_course_codes(level, program))

    completed = []
    for attempt in attempts_by_quiz.values():
        quiz_id = attempt.get('quiz_id')
        if quiz_id is None:
            continue
        quiz = quizzes_by_id.get(quiz_id)
        if not quiz:
            continue
        if level and not _is_course_allowed_for_student(quiz.get('course', ''), level, program, quiz.get('title', '')):
            continue
        completed.append({
            'quiz_id': quiz_id,
            'title': quiz.get('title'),
            'course': quiz.get('course'),
            'score': attempt.get('score', 0),
            'out_of': attempt.get('out_of', 0),
            'date': attempt.get('attempted_at'),
            'grade': attempt.get('grade', ''),
            'passed': attempt.get('passed', False),
        })

    return {'quizzes': completed}


@router.get('/{quiz_id}')
def get_quiz_details(quiz_id: int, level: str | None = None, program: str | None = None, student_id: str | None = None) -> dict[str, Any]:
    ensure_supabase_enabled()
    response = supabase.table('quizzes').select('*').eq('id', quiz_id).limit(1).execute()
    if supabase_failed(response):
        raise HTTPException(status_code=502, detail=supabase_error_message(response, 'Supabase get quiz failed'))
    if not response.data:
        raise HTTPException(status_code=404, detail='Quiz not found')

    quiz = response.data[0]
    now = datetime.now(timezone.utc)
    if _quiz_is_expired(quiz, now):
        raise HTTPException(status_code=403, detail='Quiz deadline has passed')
    if _quiz_is_locked(quiz, now):
        open_dt_str = quiz.get('open_date') or 'scheduled open time'
        raise HTTPException(status_code=403, detail=f"Quiz is locked until {open_dt_str}")

    if level:
        if not _is_course_allowed_for_student(quiz.get('course', ''), level, program, quiz.get('title', '')):
            raise HTTPException(status_code=403, detail='Quiz is not available for this student level/program')

    questions_response = supabase.table('quiz_questions').select('*').eq('quiz_id', quiz_id).execute()
    if supabase_failed(questions_response):
        raise HTTPException(status_code=502, detail=supabase_error_message(questions_response, 'Supabase get quiz questions failed'))

    raw_questions = questions_response.data or []
    import hashlib
    seed_key = f"{quiz_id}_{student_id or 'random'}"
    stable_seed = int(hashlib.md5(seed_key.encode()).hexdigest(), 16)
    rng = random.Random(stable_seed)

    questions = []
    for q in raw_questions:
        q_copy = dict(q)
        opts = q_copy.get('options')
        if isinstance(opts, list) and len(opts) > 1:
            opts_copy = list(opts)
            rng.shuffle(opts_copy)
            q_copy['options'] = opts_copy
        questions.append(q_copy)

    rng.shuffle(questions)
    return {'quiz': quiz, 'questions': questions}